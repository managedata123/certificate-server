from flask import Flask, request, send_file
from pptx import Presentation
from datetime import datetime
import uuid, os

app = Flask(__name__)

@app.route("/generate-certificates", methods=["POST"])
def generate_certificates():
    rows = request.json["rows"]
    prs = Presentation("certi_template.pptx")
    
    for row in rows:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes[0].text = row["fomat_name"]
        slide.shapes[1].text = row["upper_subject"]
        slide.shapes[2].text = row["paid_amount"]
        slide.shapes[3].text = row["period"]
        slide.shapes[4].text = "발급일: " + datetime.today().strftime("%Y-%m-%d")

    if not os.path.exists("output"):
        os.makedirs("output")
        
    filename = f"certi_{uuid.uuid4().hex}.pptx"
    path = os.path.join("output", filename)
    prs.save(path)
    
    return {"download_url": f"http://localhost:10000/download/{filename}"}

@app.route("/download/<filename>")
def download(filename):
    return send_file(os.path.join("output", filename), as_attachment=True)

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000)
